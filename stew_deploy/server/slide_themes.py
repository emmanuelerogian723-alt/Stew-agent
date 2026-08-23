"""
S.T.E.W Slide Theme Engine — 40+ professional PPTX themes.

Each theme defines:
  - A color palette (bg, panel, accent, accent_light, accent_dim, text colors)
  - A layout style (how title slides and content slides are structured)
  - A shape vocabulary (rectangles, circles, angled bars, gradient strips, etc.)
  - A category (corporate, finance, marketing, tech, healthcare, etc.)

Themes are auto-selected based on topic keywords, or can be specified by name.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import re


# ═══════════════════════════════════════════════════════════════════════════════
# THEME DEFINITIONS — 42 themes across 10 categories
# ═══════════════════════════════════════════════════════════════════════════════

THEMES = {

    # ── CORPORATE / BUSINESS (4) ──────────────────────────────────────────
    "corporate_navy": {
        "category": "corporate",
        "bg": (0x0B, 0x0F, 0x1A), "panel": (0x14, 0x1B, 0x2E),
        "accent": (0xF5, 0x9E, 0x0B), "accent_light": (0xFD, 0xD8, 0x6E), "accent_dim": (0x7A, 0x5A, 0x1E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xD7, 0xDE, 0xE8), "muted": (0x7C, 0x8A, 0xA0),
        "divider": (0x26, 0x2F, 0x45),
        "layout": "side_panel", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "corporate_blue": {
        "category": "corporate",
        "bg": (0xF8, 0xFA, 0xFC), "panel": (0x1E, 0x3A, 0x5F),
        "accent": (0x0A, 0x66, 0xB2), "accent_light": (0x3B, 0x9E, 0xD4), "accent_dim": (0xB0, 0xD4, 0xE8),
        "white": (0x1E, 0x2A, 0x38), "light": (0x4A, 0x5A, 0x6A), "muted": (0x8A, 0x9A, 0xAA),
        "divider": (0xE0, 0xE8, 0xF0),
        "layout": "top_bar", "shape": "rect", "fonts": "Calibri",
        "bg_type": "light",
    },
    "corporate_charcoal": {
        "category": "corporate",
        "bg": (0x1A, 0x1A, 0x1E), "panel": (0x2A, 0x2A, 0x32),
        "accent": (0x00, 0xCE, 0xC9), "accent_light": (0x5E, 0xFF, 0xFA), "accent_dim": (0x1A, 0x6A, 0x68),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC8, 0xCE, 0xD6), "muted": (0x6A, 0x72, 0x80),
        "divider": (0x3A, 0x3A, 0x44),
        "layout": "minimal", "shape": "rect", "fonts": "Arial",
        "bg_type": "dark",
    },
    "corporate_slate": {
        "category": "corporate",
        "bg": (0x2C, 0x3E, 0x50), "panel": (0x34, 0x49, 0x5E),
        "accent": (0xEC, 0xF0, 0xF1), "accent_light": (0xFF, 0xFF, 0xFF), "accent_dim": (0x5D, 0x6D, 0x7E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xD5, 0xDB, 0xDB), "muted": (0x85, 0x92, 0x9E),
        "divider": (0x3D, 0x55, 0x6E),
        "layout": "side_panel", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },

    # ── FINANCE / BANKING (4) ─────────────────────────────────────────────
    "banking_classic": {
        "category": "finance",
        "bg": (0x00, 0x1C, 0x3D), "panel": (0x00, 0x2B, 0x5C),
        "accent": (0xC5, 0xA5, 0x72), "accent_light": (0xE8, 0xD5, 0xA8), "accent_dim": (0x6B, 0x5A, 0x3E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC8, 0xD4, 0xE2), "muted": (0x7A, 0x8A, 0x9E),
        "divider": (0x0A, 0x3A, 0x6E),
        "layout": "bottom_bar", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "banking_trust": {
        "category": "finance",
        "bg": (0xF7, 0xF9, 0xFB), "panel": (0xE8, 0xEF, 0xF5),
        "accent": (0x00, 0x4A, 0x8F), "accent_light": (0x2E, 0x86, 0xC1), "accent_dim": (0xB0, 0xCC, 0xE0),
        "white": (0x1A, 0x2A, 0x3A), "light": (0x4A, 0x5A, 0x6A), "muted": (0x8A, 0x9A, 0xAA),
        "divider": (0xD0, 0xDE, 0xEC),
        "layout": "top_bar", "shape": "rect", "fonts": "Calibri",
        "bg_type": "light",
    },
    "finance_emerald": {
        "category": "finance",
        "bg": (0x0A, 0x1A, 0x14), "panel": (0x12, 0x2E, 0x22),
        "accent": (0x10, 0xB9, 0x81), "accent_light": (0x6E, 0xE7, 0xB6), "accent_dim": (0x1A, 0x5E, 0x44),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC8, 0xDE, 0xD4), "muted": (0x6A, 0x82, 0x76),
        "divider": (0x1A, 0x3E, 0x2E),
        "layout": "side_panel", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "finance_burgundy": {
        "category": "finance",
        "bg": (0x2A, 0x0A, 0x14), "panel": (0x3E, 0x12, 0x1E),
        "accent": (0xC8, 0x10, 0x2E), "accent_light": (0xFF, 0x6B, 0x7E), "accent_dim": (0x7E, 0x2A, 0x3A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xDE, 0xC8, 0xCE), "muted": (0x9A, 0x6A, 0x72),
        "divider": (0x4E, 0x1A, 0x2E),
        "layout": "bottom_bar", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },

    # ── INVESTMENT (4) ───────────────────────────────────────────────────
    "investment_growth": {
        "category": "investment",
        "bg": (0x0A, 0x2A, 0x1E), "panel": (0x12, 0x3E, 0x2E),
        "accent": (0x00, 0xE6, 0x76), "accent_light": (0x7B, 0xFF, 0xB8), "accent_dim": (0x2A, 0x6E, 0x4E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC8, 0xE4, 0xD4), "muted": (0x6A, 0x8A, 0x76),
        "divider": (0x1A, 0x4E, 0x3E),
        "layout": "diagonal", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "investment_gold": {
        "category": "investment",
        "bg": (0x14, 0x10, 0x00), "panel": (0x24, 0x1E, 0x08),
        "accent": (0xFF, 0xD7, 0x00), "accent_light": (0xFF, 0xEC, 0x5C), "accent_dim": (0x6E, 0x5E, 0x1A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE8, 0xDE, 0xB8), "muted": (0x9A, 0x8A, 0x5A),
        "divider": (0x34, 0x2E, 0x10),
        "layout": "side_panel", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "investment_teal": {
        "category": "investment",
        "bg": (0x00, 0x1A, 0x1E), "panel": (0x04, 0x2A, 0x30),
        "accent": (0x00, 0xB8, 0xD4), "accent_light": (0x5E, 0xE8, 0xF2), "accent_dim": (0x1A, 0x6A, 0x76),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC4, 0xE4, 0xE8), "muted": (0x5A, 0x82, 0x8A),
        "divider": (0x0A, 0x3A, 0x42),
        "layout": "minimal", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "investment_slate": {
        "category": "investment",
        "bg": (0xF4, 0xF6, 0xF8), "panel": (0xE0, 0xE8, 0xEC),
        "accent": (0x2D, 0x9B, 0x6B), "accent_light": (0x5E, 0xC4, 0x9A), "accent_dim": (0xB0, 0xD8, 0xC4),
        "white": (0x1A, 0x2A, 0x32), "light": (0x4A, 0x5A, 0x62), "muted": (0x8A, 0x9A, 0xA2),
        "divider": (0xD0, 0xDE, 0xE4),
        "layout": "top_bar", "shape": "rect", "fonts": "Calibri",
        "bg_type": "light",
    },

    # ── MARKETING (4) ─────────────────────────────────────────────────────
    "marketing_orange": {
        "category": "marketing",
        "bg": (0x1A, 0x0E, 0x00), "panel": (0x2E, 0x18, 0x04),
        "accent": (0xFF, 0x6B, 0x00), "accent_light": (0xFF, 0xA0, 0x5C), "accent_dim": (0x7A, 0x3A, 0x10),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE8, 0xD4, 0xC4), "muted": (0x9A, 0x7A, 0x6A),
        "divider": (0x3E, 0x24, 0x0A),
        "layout": "diagonal", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "marketing_magenta": {
        "category": "marketing",
        "bg": (0x1A, 0x00, 0x1E), "panel": (0x2E, 0x08, 0x34),
        "accent": (0xE9, 0x1E, 0x63), "accent_light": (0xFF, 0x6B, 0xA0), "accent_dim": (0x7E, 0x2A, 0x4E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE8, 0xC4, 0xD4), "muted": (0x9A, 0x6A, 0x7A),
        "divider": (0x3E, 0x14, 0x42),
        "layout": "bottom_bar", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "marketing_coral": {
        "category": "marketing",
        "bg": (0xFF, 0xF5, 0xF0), "panel": (0xFF, 0xE8, 0xDE),
        "accent": (0xFF, 0x5C, 0x5C), "accent_light": (0xFF, 0x8E, 0x8E), "accent_dim": (0xFF, 0xC4, 0xC4),
        "white": (0x2A, 0x1A, 0x1A), "light": (0x5A, 0x4A, 0x4A), "muted": (0x9A, 0x8A, 0x8A),
        "divider": (0xFF, 0xDE, 0xD0),
        "layout": "top_bar", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "marketing_sunset": {
        "category": "marketing",
        "bg": (0x2A, 0x0A, 0x2A), "panel": (0x3E, 0x12, 0x3E),
        "accent": (0xFF, 0x8C, 0x42), "accent_light": (0xFF, 0xC4, 0x8C), "accent_dim": (0x7E, 0x4A, 0x2E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE8, 0xD4, 0xE4), "muted": (0x9A, 0x7A, 0x8E),
        "divider": (0x4E, 0x1E, 0x4E),
        "layout": "side_panel", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },

    # ── TECHNOLOGY / STARTUP (4) ─────────────────────────────────────────
    "tech_neon": {
        "category": "tech",
        "bg": (0x00, 0x00, 0x00), "panel": (0x0A, 0x0A, 0x12),
        "accent": (0x00, 0xFF, 0x9F), "accent_light": (0x5E, 0xFF, 0xC8), "accent_dim": (0x1A, 0x6E, 0x4E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC8, 0xDE, 0xD4), "muted": (0x5A, 0x72, 0x6A),
        "divider": (0x1A, 0x1A, 0x24),
        "layout": "minimal", "shape": "rect", "fonts": "Consolas",
        "bg_type": "dark",
    },
    "tech_purple": {
        "category": "tech",
        "bg": (0x14, 0x0A, 0x2E), "panel": (0x24, 0x12, 0x42),
        "accent": (0x8B, 0x5C, 0xF6), "accent_light": (0xC4, 0xA8, 0xFF), "accent_dim": (0x5A, 0x3E, 0x9A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xD8, 0xCE, 0xE8), "muted": (0x8A, 0x7A, 0xA2),
        "divider": (0x34, 0x1E, 0x52),
        "layout": "diagonal", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "tech_cyan": {
        "category": "tech",
        "bg": (0x00, 0x14, 0x1E), "panel": (0x04, 0x24, 0x34),
        "accent": (0x00, 0xD4, 0xFF), "accent_light": (0x5E, 0xE8, 0xFF), "accent_dim": (0x1A, 0x6A, 0x84),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC4, 0xE4, 0xEC), "muted": (0x5A, 0x82, 0x8E),
        "divider": (0x0A, 0x34, 0x42),
        "layout": "side_panel", "shape": "rect", "fonts": "Consolas",
        "bg_type": "dark",
    },
    "tech_startup": {
        "category": "tech",
        "bg": (0x0A, 0x0E, 0x1E), "panel": (0x14, 0x1E, 0x34),
        "accent": (0x00, 0x9F, 0xFF), "accent_light": (0x5C, 0xC4, 0xFF), "accent_dim": (0x2A, 0x5A, 0x8E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC4, 0xD4, 0xE4), "muted": (0x5A, 0x6E, 0x8A),
        "divider": (0x1E, 0x2E, 0x4E),
        "layout": "bottom_bar", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },

    # ── EDUCATION (4) ────────────────────────────────────────────────────
    "education_blue": {
        "category": "education",
        "bg": (0xF0, 0xF4, 0xF8), "panel": (0xE0, 0xEC, 0xF4),
        "accent": (0x00, 0x66, 0xB2), "accent_light": (0x3B, 0x9E, 0xD4), "accent_dim": (0xB0, 0xD4, 0xE8),
        "white": (0x1A, 0x2A, 0x38), "light": (0x4A, 0x5A, 0x6A), "muted": (0x8A, 0x9A, 0xAA),
        "divider": (0xD0, 0xDE, 0xEC),
        "layout": "top_bar", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "education_green": {
        "category": "education",
        "bg": (0xF4, 0xF8, 0xF0), "panel": (0xE4, 0xEC, 0xDC),
        "accent": (0x2D, 0x9B, 0x3D), "accent_light": (0x5E, 0xC4, 0x6E), "accent_dim": (0xB0, 0xD8, 0xB4),
        "white": (0x1A, 0x2A, 0x1A), "light": (0x4A, 0x5A, 0x4A), "muted": (0x8A, 0x9A, 0x8A),
        "divider": (0xD0, 0xDE, 0xCE),
        "layout": "top_bar", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "education_amber": {
        "category": "education",
        "bg": (0xFF, 0xFB, 0xF0), "panel": (0xFF, 0xF0, 0xD0),
        "accent": (0xF5, 0x9E, 0x0B), "accent_light": (0xFD, 0xD8, 0x6E), "accent_dim": (0xFF, 0xE4, 0xA8),
        "white": (0x2A, 0x24, 0x10), "light": (0x5A, 0x4E, 0x3A), "muted": (0x9A, 0x8E, 0x7A),
        "divider": (0xFF, 0xEC, 0xC4),
        "layout": "side_panel", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "education_indigo": {
        "category": "education",
        "bg": (0x0E, 0x0A, 0x2E), "panel": (0x1E, 0x14, 0x42),
        "accent": (0x7C, 0x5C, 0xF6), "accent_light": (0xB4, 0x9C, 0xFF), "accent_dim": (0x4A, 0x3A, 0x9A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xD4, 0xCE, 0xE8), "muted": (0x8A, 0x7E, 0xA2),
        "divider": (0x2A, 0x1E, 0x52),
        "layout": "minimal", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },

    # ── HEALTHCARE (4) ───────────────────────────────────────────────────
    "healthcare_teal": {
        "category": "healthcare",
        "bg": (0xF0, 0xF8, 0xFA), "panel": (0xDC, 0xF0, 0xF4),
        "accent": (0x00, 0x9F, 0xB2), "accent_light": (0x5E, 0xC4, 0xCE), "accent_dim": (0xB0, 0xDE, 0xE4),
        "white": (0x1A, 0x2A, 0x2E), "light": (0x4A, 0x5A, 0x5E), "muted": (0x8A, 0x9A, 0x9E),
        "divider": (0xD0, 0xE4, 0xE8),
        "layout": "top_bar", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "healthcare_mint": {
        "category": "healthcare",
        "bg": (0xF4, 0xFA, 0xF6), "panel": (0xE0, 0xF0, 0xE8),
        "accent": (0x00, 0xB3, 0x7A), "accent_light": (0x5E, 0xD4, 0xAE), "accent_dim": (0xB0, 0xDE, 0xCE),
        "white": (0x1A, 0x2A, 0x24), "light": (0x4A, 0x5A, 0x52), "muted": (0x8A, 0x9A, 0x92),
        "divider": (0xD0, 0xE4, 0xDC),
        "layout": "side_panel", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "healthcare_blue": {
        "category": "healthcare",
        "bg": (0xF8, 0xFB, 0xFF), "panel": (0xE4, 0xEC, 0xF8),
        "accent": (0x00, 0x6B, 0xD4), "accent_light": (0x5E, 0xA4, 0xE4), "accent_dim": (0xB0, 0xCC, 0xE4),
        "white": (0x1A, 0x24, 0x2E), "light": (0x4A, 0x5A, 0x62), "muted": (0x8A, 0x9A, 0x9E),
        "divider": (0xD4, 0xE0, 0xEC),
        "layout": "top_bar", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "healthcare_clean": {
        "category": "healthcare",
        "bg": (0xFF, 0xFF, 0xFF), "panel": (0xF4, 0xF8, 0xFC),
        "accent": (0x2D, 0xA8, 0x96), "accent_light": (0x5E, 0xCE, 0xB8), "accent_dim": (0xC4, 0xE8, 0xE0),
        "white": (0x1A, 0x1E, 0x22), "light": (0x4A, 0x4E, 0x52), "muted": (0x8A, 0x8E, 0x92),
        "divider": (0xE4, 0xE8, 0xEC),
        "layout": "minimal", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },

    # ── CREATIVE / AGENCY (4) ────────────────────────────────────────────
    "creative_purple": {
        "category": "creative",
        "bg": (0x14, 0x08, 0x1E), "panel": (0x24, 0x10, 0x34),
        "accent": (0xD0, 0x3A, 0xFF), "accent_light": (0xE8, 0x8C, 0xFF), "accent_dim": (0x7A, 0x2A, 0x9A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE0, 0xCE, 0xE8), "muted": (0x9A, 0x7A, 0xA2),
        "divider": (0x34, 0x1E, 0x44),
        "layout": "diagonal", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "creative_pink": {
        "category": "creative",
        "bg": (0xFF, 0xF0, 0xF8), "panel": (0xFF, 0xE0, 0xEC),
        "accent": (0xE9, 0x1E, 0x8E), "accent_light": (0xFF, 0x6B, 0xC4), "accent_dim": (0xFF, 0xC0, 0xDE),
        "white": (0x2A, 0x1A, 0x24), "light": (0x5A, 0x4A, 0x52), "muted": (0x9A, 0x8A, 0x92),
        "divider": (0xFF, 0xD4, 0xE4),
        "layout": "top_bar", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "creative_lime": {
        "category": "creative",
        "bg": (0x0A, 0x1A, 0x0A), "panel": (0x12, 0x2E, 0x12),
        "accent": (0xB4, 0xFF, 0x00), "accent_light": (0xD4, 0xFF, 0x5E), "accent_dim": (0x5E, 0x8E, 0x1A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xCE, 0xDE, 0xCE), "muted": (0x6E, 0x82, 0x6E),
        "divider": (0x1A, 0x3E, 0x1A),
        "layout": "minimal", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "creative_bold": {
        "category": "creative",
        "bg": (0x00, 0x00, 0x00), "panel": (0x14, 0x14, 0x14),
        "accent": (0xFF, 0x00, 0x55), "accent_light": (0xFF, 0x5E, 0x9A), "accent_dim": (0x7A, 0x2A, 0x44),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE4, 0xE4, 0xE4), "muted": (0x8A, 0x8A, 0x8A),
        "divider": (0x24, 0x24, 0x24),
        "layout": "bottom_bar", "shape": "rect", "fonts": "Arial",
        "bg_type": "dark",
    },

    # ── MINIMAL / CLEAN (4) ──────────────────────────────────────────────
    "minimal_white": {
        "category": "minimal",
        "bg": (0xFF, 0xFF, 0xFF), "panel": (0xF8, 0xFA, 0xFC),
        "accent": (0x2D, 0x6B, 0xFF), "accent_light": (0x7A, 0xA8, 0xFF), "accent_dim": (0xC4, 0xD8, 0xFF),
        "white": (0x1A, 0x1E, 0x24), "light": (0x4A, 0x4E, 0x52), "muted": (0x8A, 0x8E, 0x92),
        "divider": (0xE4, 0xE8, 0xEC),
        "layout": "minimal", "shape": "rect", "fonts": "Calibri",
        "bg_type": "light",
    },
    "minimal_gray": {
        "category": "minimal",
        "bg": (0xF4, 0xF4, 0xF4), "panel": (0xEC, 0xEC, 0xEC),
        "accent": (0x5A, 0x5A, 0x5A), "accent_light": (0x8A, 0x8A, 0x8A), "accent_dim": (0xC4, 0xC4, 0xC4),
        "white": (0x1A, 0x1A, 0x1A), "light": (0x4A, 0x4A, 0x4A), "muted": (0x8A, 0x8A, 0x8A),
        "divider": (0xDE, 0xDE, 0xDE),
        "layout": "minimal", "shape": "rect", "fonts": "Calibri",
        "bg_type": "light",
    },
    "minimal_cream": {
        "category": "minimal",
        "bg": (0xFA, 0xF6, 0xF0), "panel": (0xF2, 0xEC, 0xE4),
        "accent": (0x8E, 0x6E, 0x4A), "accent_light": (0xC4, 0xA4, 0x7A), "accent_dim": (0xDE, 0xCE, 0xB4),
        "white": (0x2A, 0x24, 0x1E), "light": (0x5A, 0x4E, 0x42), "muted": (0x9A, 0x8E, 0x82),
        "divider": (0xE8, 0xE0, 0xD4),
        "layout": "side_panel", "shape": "rect", "fonts": "Georgia",
        "bg_type": "light",
    },
    "minimal_snow": {
        "category": "minimal",
        "bg": (0xFC, 0xFE, 0xFF), "panel": (0xF0, 0xF4, 0xF8),
        "accent": (0x00, 0xB3, 0xD4), "accent_light": (0x5E, 0xCE, 0xE4), "accent_dim": (0xC4, 0xE4, 0xEC),
        "white": (0x1A, 0x1E, 0x22), "light": (0x4A, 0x4E, 0x52), "muted": (0x8A, 0x8E, 0x92),
        "divider": (0xE4, 0xEC, 0xF0),
        "layout": "minimal", "shape": "rect", "fonts": "Calibri",
        "bg_type": "light",
    },

    # ── NATURE / ECO (4) ─────────────────────────────────────────────────
    "nature_forest": {
        "category": "nature",
        "bg": (0x0A, 0x1E, 0x0E), "panel": (0x12, 0x2E, 0x16),
        "accent": (0x4C, 0xAF, 0x50), "accent_light": (0x8B, 0xD6, 0x8E), "accent_dim": (0x2A, 0x6E, 0x2E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xCE, 0xDE, 0xCE), "muted": (0x6E, 0x82, 0x6E),
        "divider": (0x1A, 0x3E, 0x1E),
        "layout": "side_panel", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },
    "nature_earth": {
        "category": "nature",
        "bg": (0x1E, 0x14, 0x08), "panel": (0x2E, 0x1E, 0x10),
        "accent": (0x8D, 0x6E, 0x63), "accent_light": (0xC4, 0xA4, 0x94), "accent_dim": (0x5E, 0x4A, 0x3E),
        "white": (0xFF, 0xF8, 0xF4), "light": (0xE4, 0xCE, 0xC4), "muted": (0x9A, 0x82, 0x72),
        "divider": (0x3E, 0x2A, 0x1A),
        "layout": "bottom_bar", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "nature_sage": {
        "category": "nature",
        "bg": (0xF0, 0xF4, 0xEC), "panel": (0xE0, 0xE8, 0xD4),
        "accent": (0x6B, 0x8E, 0x4A), "accent_light": (0x9A, 0xC4, 0x72), "accent_dim": (0xCE, 0xDE, 0xB4),
        "white": (0x1E, 0x24, 0x14), "light": (0x4A, 0x4E, 0x3A), "muted": (0x82, 0x8A, 0x72),
        "divider": (0xD0, 0xDC, 0xC4),
        "layout": "top_bar", "shape": "rounded", "fonts": "Calibri",
        "bg_type": "light",
    },
    "nature_olive": {
        "category": "nature",
        "bg": (0x1A, 0x1E, 0x08), "panel": (0x2A, 0x2E, 0x12),
        "accent": (0xB4, 0xC4, 0x2D), "accent_light": (0xD4, 0xDE, 0x72), "accent_dim": (0x6E, 0x7E, 0x2A),
        "white": (0xFF, 0xFE, 0xF4), "light": (0xDE, 0xDE, 0xCE), "muted": (0x8A, 0x8A, 0x72),
        "divider": (0x3A, 0x3E, 0x1A),
        "layout": "minimal", "shape": "rect", "fonts": "Segoe UI",
        "bg_type": "dark",
    },

    # ── AFRICAN / HERITAGE (4) ───────────────────────────────────────────
    "african_gold": {
        "category": "african",
        "bg": (0x1A, 0x0E, 0x00), "panel": (0x2E, 0x18, 0x04),
        "accent": (0xFF, 0xB3, 0x00), "accent_light": (0xFF, 0xD4, 0x5C), "accent_dim": (0x7E, 0x5A, 0x1A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE4, 0xD4, 0xC4), "muted": (0x9A, 0x82, 0x6E),
        "divider": (0x3E, 0x24, 0x08),
        "layout": "diagonal", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "african_sunset": {
        "category": "african",
        "bg": (0x2A, 0x10, 0x00), "panel": (0x3E, 0x18, 0x04),
        "accent": (0xFF, 0x6B, 0x1A), "accent_light": (0xFF, 0xA4, 0x5C), "accent_dim": (0x7E, 0x3A, 0x1E),
        "white": (0xFF, 0xF8, 0xF0), "light": (0xE4, 0xCE, 0xBE), "muted": (0x9A, 0x7E, 0x6E),
        "divider": (0x4E, 0x24, 0x08),
        "layout": "bottom_bar", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "african_red": {
        "category": "african",
        "bg": (0x1E, 0x04, 0x00), "panel": (0x2E, 0x08, 0x04),
        "accent": (0xC8, 0x2A, 0x1A), "accent_light": (0xFF, 0x6E, 0x5C), "accent_dim": (0x7E, 0x2A, 0x22),
        "white": (0xFF, 0xF8, 0xF4), "light": (0xE4, 0xCE, 0xC4), "muted": (0x9A, 0x7E, 0x72),
        "divider": (0x3E, 0x0E, 0x08),
        "layout": "side_panel", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "african_kente": {
        "category": "african",
        "bg": (0x0A, 0x0A, 0x00), "panel": (0x1A, 0x1A, 0x04),
        "accent": (0xFF, 0xC4, 0x00), "accent_light": (0xFF, 0xE4, 0x5C), "accent_dim": (0x7A, 0x6E, 0x1E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE4, 0xE4, 0xC4), "muted": (0x8E, 0x8A, 0x6A),
        "divider": (0x2A, 0x2A, 0x08),
        "layout": "diagonal", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },

    # ── LUXURY / PREMIUM (2) ─────────────────────────────────────────────
    "luxury_black_gold": {
        "category": "luxury",
        "bg": (0x00, 0x00, 0x00), "panel": (0x0A, 0x0A, 0x0A),
        "accent": (0xD4, 0xAF, 0x37), "accent_light": (0xF0, 0xDE, 0x8C), "accent_dim": (0x6E, 0x5E, 0x2A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE4, 0xDE, 0xC4), "muted": (0x8A, 0x82, 0x6E),
        "divider": (0x1A, 0x1A, 0x1A),
        "layout": "side_panel", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "luxury_charcoal_silver": {
        "category": "luxury",
        "bg": (0x1C, 0x1C, 0x1E), "panel": (0x28, 0x28, 0x2A),
        "accent": (0xC0, 0xC0, 0xC0), "accent_light": (0xE4, 0xE4, 0xE4), "accent_dim": (0x6E, 0x6E, 0x6E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xD4, 0xD4, 0xD6), "muted": (0x7A, 0x7A, 0x7C),
        "divider": (0x34, 0x34, 0x36),
        "layout": "minimal", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    # ── CHURCH / RELIGION (4) ────────────────────────────────────────────
    "church_royal": {
        "category": "church",
        "bg": (0x0A, 0x04, 0x1E), "panel": (0x12, 0x0A, 0x34),
        "accent": (0xD4, 0xAF, 0x37), "accent_light": (0xF0, 0xDE, 0x8C), "accent_dim": (0x6E, 0x5E, 0x2A),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xE4, 0xDE, 0xC4), "muted": (0x8A, 0x82, 0x6E),
        "divider": (0x1A, 0x12, 0x3E),
        "layout": "side_panel", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "church_warm": {
        "category": "church",
        "bg": (0x1E, 0x10, 0x08), "panel": (0x2E, 0x1C, 0x10),
        "accent": (0xE8, 0x8C, 0x2A), "accent_light": (0xFF, 0xC4, 0x72), "accent_dim": (0x7E, 0x4E, 0x1E),
        "white": (0xFF, 0xF8, 0xF0), "light": (0xE4, 0xCE, 0xBE), "muted": (0x9A, 0x82, 0x6E),
        "divider": (0x3E, 0x2A, 0x1A),
        "layout": "bottom_bar", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },
    "church_divine": {
        "category": "church",
        "bg": (0xF8, 0xF4, 0xF0), "panel": (0xF0, 0xE8, 0xE0),
        "accent": (0x8B, 0x6F, 0x3E), "accent_light": (0xC4, 0xA4, 0x6E), "accent_dim": (0xDE, 0xCE, 0xAE),
        "white": (0x2A, 0x20, 0x14), "light": (0x5A, 0x4E, 0x3A), "muted": (0x9A, 0x8E, 0x7A),
        "divider": (0xE4, 0xDE, 0xCE),
        "layout": "top_bar", "shape": "rounded", "fonts": "Georgia",
        "bg_type": "light",
    },
    "church_celestial": {
        "category": "church",
        "bg": (0x04, 0x0E, 0x1E), "panel": (0x0A, 0x1A, 0x34),
        "accent": (0x5C, 0xB4, 0xFF), "accent_light": (0xA4, 0xD4, 0xFF), "accent_dim": (0x2A, 0x5A, 0x8E),
        "white": (0xFF, 0xFF, 0xFF), "light": (0xC4, 0xD8, 0xE4), "muted": (0x6A, 0x82, 0x92),
        "divider": (0x0A, 0x2A, 0x3E),
        "layout": "minimal", "shape": "rect", "fonts": "Georgia",
        "bg_type": "dark",
    },

}

# ═══════════════════════════════════════════════════════════════════════════════
# AUTO-THEME DETECTION — pick best theme from topic keywords
# ═══════════════════════════════════════════════════════════════════════════════

THEME_KEYWORDS = {
    "finance": ["bank", "banking", "financial", "finance", "loan", "credit", "mortgage",
                "interest rate", "fintech", "payment", "paystack", "transaction"],
    "investment": ["invest", "investment", "portfolio", "stock", "dividend", "ROI",
                   "asset", "capital", "venture", "fund", "equity", "trading"],
    "marketing": ["market", "marketing", "campaign", "brand", "advertis", "social media",
                  "content strategy", "SEO", "promotion", "growth hack"],
    "tech": ["tech", "technology", "software", "app", "startup", "AI", "artificial intelligence",
             "machine learning", "blockchain", "crypto", "SaaS", "platform", "API", "developer"],
    "education": ["school", "education", "lecture", "course", "training", "learn",
                  "teach", "student", "curriculum", "academic", "university", "quiz"],
    "healthcare": ["health", "healthcare", "medical", "hospital", "clinic", "patient",
                   "doctor", "wellness", "pharma", "disease", "treatment", "nursing"],
    "creative": ["design", "creative", "agency", "branding", "art", "studio",
                 "portfolio", "visual", "aesthetic", "inspiration"],
    "nature": ["environment", "eco", "green", "sustainab", "climate", "nature",
               "forest", "renewable", "carbon", "conservation", "organic"],
    "african": ["africa", "nigeria", "nigerian", "lagos", "abuja", "naira", "afro",
               "kente", "heritage", "ghanian", "kenyan", "ghana", "kenya"],
    "luxury": ["luxury", "premium", "exclusive", "elegant", "sophisticated",
               "high-end", "boutique", "vip", "platinum", "gold member"],
    "corporate": ["corporate", "business", "company", "enterprise", "organization",
                  "strategy", "report", "annual", "quarterly", "executive"],
    "minimal": ["minimal", "simple", "clean", "basic", "overview", "summary",
                "brief", "introduction"],
    "church": ["church", "fundrais", "giving", "donation", "tithe", "offering",
              "ministry", "pastor", "congregation", "faith", "gospel", "christian",
              "prayer", "worship", "charity", "religious", "spiritual", "temple"],
}


# High-priority keywords that strongly indicate a specific category
# (weight 3) vs normal keywords (weight 1)
HIGH_PRIORITY_KEYWORDS = {
    "finance": ["bank", "banking", "fintech", "mortgage", "interest rate"],
    "investment": ["invest", "investment", "portfolio", "dividend", "trading"],
    "tech": ["AI", "artificial intelligence", "machine learning", "blockchain", "crypto"],
    "healthcare": ["hospital", "clinic", "patient", "doctor", "pharma"],
    "african": ["africa", "nigeria", "nigerian", "naira", "kente"],
    "luxury": ["luxury", "platinum", "vip"],
    "church": ["church", "fundrais", "tithe", "offering", "ministry", "pastor", "gospel"],
}


def auto_select_theme(topic: str) -> str:
    """Pick the best theme name based on keywords in the topic/title."""
    topic_lower = topic.lower()

    # Score each category by keyword matches (with weighted priorities)
    scores = {}
    for category, keywords in THEME_KEYWORDS.items():
        score = sum(1 for kw in keywords if kw in topic_lower)
        # Apply high-priority weight boost
        high_kw = HIGH_PRIORITY_KEYWORDS.get(category, [])
        score += sum(2 for kw in high_kw if kw in topic_lower)  # +2 per high-priority match
        if score > 0:
            scores[category] = score

    if not scores:
        # Default: corporate navy for general business presentations
        return "corporate_navy"

    best_category = max(scores, key=scores.get)

    # Pick the first theme in the best category
    for name, theme in THEMES.items():
        if theme["category"] == best_category:
            return name

    return "corporate_navy"


def get_theme(theme_name: str = None, topic: str = "") -> dict:
    """Get a theme by name, or auto-select based on topic."""
    if theme_name and theme_name in THEMES:
        return THEMES[theme_name]
    return THEMES[auto_select_theme(topic)]


def list_themes() -> list:
    """Return all available theme names with their categories."""
    return [
        {"name": name, "category": theme["category"], "bg_type": theme["bg_type"]}
        for name, theme in THEMES.items()
    ]


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE RENDERER — renders slides using a theme
# ═══════════════════════════════════════════════════════════════════════════════

def _c(rgb_tuple):
    """Convert (r, g, b) tuple to RGBColor."""
    return RGBColor(*rgb_tuple)


# ═══════════════════════════════════════════════════════════════════════════════
# HERO IMAGE FETCHING — real photography for Canva-style title/closing slides
# ═══════════════════════════════════════════════════════════════════════════════

CATEGORY_IMAGE_STYLE = {
    "corporate": "modern corporate office skyline, glass towers, professional, cinematic lighting",
    "finance": "modern bank headquarters, financial district skyline at dusk, cinematic",
    "investment": "stock market growth chart, city skyline, golden hour, cinematic finance",
    "marketing": "vibrant creative agency workspace, colorful abstract shapes, dynamic energy",
    "tech": "futuristic technology data center, neon blue lights, servers, sci-fi atmosphere",
    "education": "bright modern classroom, students learning, warm natural light",
    "healthcare": "modern hospital interior, clean bright medical facility, calm atmosphere",
    "creative": "vibrant art studio, colorful paint splashes, creative energy, bold colors",
    "minimal": "minimalist white studio space, soft natural light, clean aesthetic",
    "nature": "lush green forest, sunlight through trees, serene natural landscape",
    "african": "vibrant African market scene, warm golden light, rich cultural colors",
    "luxury": "elegant black marble interior, gold accents, luxury architecture",
    "church": "warm cathedral interior, sunlight through stained glass windows, cinematic",
}


def _fetch_hero_image(topic: str, category: str, seed: int = None, timeout: int = 18, retries: int = 3):
    """Fetch a real AI-generated background photo for hero slides via pollinations.ai.
    Returns raw JPEG bytes, or None if the fetch fails (caller must handle gracefully).
    Retries with backoff on 429 (rate limit) since consecutive calls (hero + closing
    image) can trigger pollinations' free-tier rate limiting."""
    import time as _time
    try:
        import httpx
        import urllib.parse
        import random as _random

        style = CATEGORY_IMAGE_STYLE.get(category, "professional abstract background, cinematic lighting")
        # Keep topic short to avoid overly literal/cluttered renders
        short_topic = re.sub(r'[^a-zA-Z0-9 ]', '', topic)[:60]
        prompt = f"{style}, related to {short_topic}, no text, no words, photography, 4k"
        encoded = urllib.parse.quote(prompt, safe='')
        seed = seed or _random.randint(1, 999999)
        url = f"https://image.pollinations.ai/prompt/{encoded}?width=1280&height=720&model=flux&nologo=true&seed={seed}"

        for attempt in range(retries):
            try:
                with httpx.Client(timeout=timeout) as client:
                    resp = client.get(url)
                    if resp.status_code == 200 and len(resp.content) > 2000:
                        return resp.content
                    elif resp.status_code == 429:
                        _time.sleep(3 + attempt * 2)  # backoff: 3s, 5s, 7s
                        continue
                    else:
                        break
            except Exception:
                _time.sleep(2)
                continue
    except Exception:
        pass
    return None


def _set_shape_alpha(shape, alpha_percent: int):
    """Apply real alpha transparency to a shape's solid fill via raw OOXML.
    alpha_percent: 0-100, where 100 = fully opaque, 0 = fully transparent."""
    try:
        val = str(int(alpha_percent * 1000))
        sp = shape.fill.fore_color._xFill
        srgbClr = sp.find(qn('a:srgbClr'))
        if srgbClr is not None:
            existing = srgbClr.find(qn('a:alpha'))
            if existing is not None:
                srgbClr.remove(existing)
            alpha_el = srgbClr.makeelement(qn('a:alpha'), {'val': val})
            srgbClr.append(alpha_el)
    except Exception:
        pass


def render_pptx(prs, slides, title, theme_name=None, total_slides=None, hero_image=None, closing_image=None):
    """Render slides into a Presentation object using the specified theme.
    Enhanced with decorative shapes, alternating content layouts, stat blocks,
    and gradient-like effects for world-class visual quality.
    """
    theme = get_theme(theme_name, title)

    # Unpack colors
    BG = _c(theme["bg"])
    BG_PANEL = _c(theme["panel"])
    ACCENT = _c(theme["accent"])
    ACCENT_LIGHT = _c(theme["accent_light"])
    ACCENT_DIM = _c(theme["accent_dim"])
    WHITE = _c(theme["white"])
    LIGHT = _c(theme["light"])
    MUTED = _c(theme["muted"])
    DIVIDER = _c(theme["divider"])
    FONT = theme.get("fonts", "Segoe UI")
    LAYOUT = theme.get("layout", "side_panel")
    SHAPE_TYPE = theme.get("shape", "rect")
    BG_DARK = theme.get("bg_type", "dark") == "dark"

    BLANK = prs.slide_layouts[6]
    SLIDE_W = 13.333
    SLIDE_H = 7.5
    if total_slides is None:
        total_slides = len(slides)

    shape_enum = MSO_SHAPE.ROUNDED_RECTANGLE if SHAPE_TYPE == "rounded" else MSO_SHAPE.RECTANGLE

    def _add_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

    def _shape(slide, left, top, width, height, color, shape_type=None):
        st = shape_type or shape_enum
        s = slide.shapes.add_shape(st, Inches(left), Inches(top), Inches(width), Inches(height))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        s.shadow.inherit = False
        return s

    def _rect(slide, left, top, width, height, color):
        return _shape(slide, left, top, width, height, color, MSO_SHAPE.RECTANGLE)

    def _rrect(slide, left, top, width, height, color):
        return _shape(slide, left, top, width, height, color, MSO_SHAPE.ROUNDED_RECTANGLE)

    def _oval(slide, left, top, size, color):
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        s.shadow.inherit = False
        return s

    def _tri(slide, left, top, width, height, color):
        s = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        s.shadow.inherit = False
        return s

    def _chevron(slide, left, top, width, height, color):
        s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(top), Inches(width), Inches(height))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        s.shadow.inherit = False
        return s

    def _text(slide, left, top, width, height, text, size, color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              font=FONT, spacing=1.0):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        para = tf.paragraphs[0]
        clean = _sanitize(text)
        para.text = clean
        para.alignment = align
        para.line_spacing = spacing
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
        return tf

    def _eyebrow(slide, left, top, text):
        _text(slide, left, top, 8.0, 0.4, text.upper(), 12.5, ACCENT, bold=True, font=FONT)

    def _page_footer(slide, index):
        _rect(slide, 0, SLIDE_H - 0.06, SLIDE_W, 0.06, ACCENT)
        _text(slide, SLIDE_W - 2.2, SLIDE_H - 0.55, 1.9, 0.4,
              f"{index+1:02d}  /  {total_slides:02d}", 11, MUTED, align=PP_ALIGN.RIGHT)
        _text(slide, 0.6, SLIDE_H - 0.55, 4.0, 0.4, "S.T.E.W", 11, MUTED, bold=True)

    def _bullets(slide, left, top, width, bullets, font_size=17.5, line_gap=0.62):
        y = top
        for bullet in bullets:
            clean = _sanitize(bullet.strip().lstrip("-").lstrip("*").lstrip("\u2022").strip())
            if not clean:
                continue
            _rect(slide, left, y + 0.12, 0.14, 0.14, ACCENT)
            box = slide.shapes.add_textbox(Inches(left + 0.35), Inches(y), Inches(width - 0.35), Inches(line_gap + 0.3))
            tf = box.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = clean
            para.line_spacing = 1.05
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(font_size)
            run.font.color.rgb = LIGHT
            run.font.name = FONT
            y += line_gap
        return y

    def _decorative_shapes(slide, variant=0):
        """Add decorative geometric shapes based on variant for visual richness."""
        if variant == 0:
            # Large faint circles top-right + bottom-left
            _oval(slide, SLIDE_W - 1.8, -0.6, 2.4, ACCENT_DIM)
            _oval(slide, -0.8, SLIDE_H - 1.8, 2.0, ACCENT_DIM)
        elif variant == 1:
            # Triangle accent + circle
            _tri(slide, SLIDE_W - 2.5, 0, 2.5, 2.5, ACCENT_DIM)
            _oval(slide, -0.5, SLIDE_H - 2.0, 1.8, ACCENT_DIM)
        elif variant == 2:
            # Chevron strip on the right
            _chevron(slide, SLIDE_W - 1.2, 1.5, 1.0, 4.5, ACCENT_DIM)
            _oval(slide, 10.5, -0.5, 1.5, ACCENT_DIM)
        elif variant == 3:
            # Two circles + small triangle
            _oval(slide, SLIDE_W - 2.0, SLIDE_H - 2.5, 2.5, ACCENT_DIM)
            _oval(slide, -0.6, -0.6, 1.8, ACCENT_DIM)
            _tri(slide, 11.0, 0.2, 1.5, 1.5, ACCENT_DIM)
        else:
            # Parallelogram accent
            _shape(slide, SLIDE_W - 3, SLIDE_H - 2.5, 3.5, 2.5, ACCENT_DIM, MSO_SHAPE.PARALLELOGRAM)
            _oval(slide, -0.4, -0.4, 1.5, ACCENT_DIM)

    def _gradient_strip(slide, left, top, width, height):
        """Simulate a gradient bar using overlapping shapes with decreasing opacity colors."""
        # Use 3 segments from accent to accent_dim for a gradient-like effect
        seg = width / 3
        _rect(slide, left, top, seg, height, ACCENT)
        _rect(slide, left + seg, top, seg, height, ACCENT_LIGHT)
        _rect(slide, left + seg * 2, top, seg, height, ACCENT_DIM)

    def _stat_block(slide, left, top, width, number, label):
        """A highlighted stat/number block for data-heavy slides."""
        _rrect(slide, left, top, width, 2.0, BG_PANEL)
        _rect(slide, left, top, 0.08, 2.0, ACCENT)
        _text(slide, left + 0.3, top + 0.25, width - 0.5, 1.0, number, 40, ACCENT_LIGHT, bold=True)
        _text(slide, left + 0.3, top + 1.25, width - 0.5, 0.6, label, 13, MUTED)

    def _card(slide, left, top, width, height, title_text, body_text):
        """A content card with accent top border."""
        _rrect(slide, left, top, width, height, BG_PANEL)
        _rect(slide, left, top, width, 0.06, ACCENT)
        _text(slide, left + 0.25, top + 0.2, width - 0.4, 0.5, title_text, 16, ACCENT_LIGHT, bold=True)
        _text(slide, left + 0.25, top + 0.75, width - 0.4, height - 0.9, body_text, 13, LIGHT, spacing=1.1)

    # ── Render each slide ──
    for i, slide_data in enumerate(slides):
        slide_title = slide_data.get("title", f"Slide {i+1}")
        slide_content = slide_data.get("content", "")
        bullets = [b.strip() for b in slide_content.split("\n") if b.strip()]

        slide = prs.slides.add_slide(BLANK)
        _add_bg(slide)

        if i == 0:
            # ══ TITLE SLIDE ══ — Canva-style hero image when available, themed flat design otherwise
            if hero_image:
                import io as _io
                try:
                    slide.shapes.add_picture(_io.BytesIO(hero_image), Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
                except Exception:
                    hero_image = None

            if hero_image and LAYOUT == "side_panel":
                overlay = _rect(slide, 0, 0, 5.6, SLIDE_H, BG_PANEL)
                _set_shape_alpha(overlay, 82)
                _rect(slide, 5.6, 0, 0.05, SLIDE_H, ACCENT)
                _eyebrow(slide, 0.5, 0.6, "PRESENTATION")
                _text(slide, 0.5, 1.4, 4.8, 2.2, slide_title, 34, WHITE, bold=True, spacing=1.0)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _rect(slide, 0.55, 3.5, 0.7, 0.045, ACCENT)
                _text(slide, 0.5, 3.7, 4.8, 1.6, subtitle, 15, ACCENT_LIGHT, spacing=1.15)
                _text(slide, 0.5, SLIDE_H - 0.9, 4.8, 0.5, "Generated by S.T.E.W Agent", 11, MUTED)
                _text(slide, 0.5, SLIDE_H - 1.3, 4.8, 0.4, _date_str(), 11, MUTED)

            elif hero_image and LAYOUT == "top_bar":
                overlay = _rect(slide, 0, SLIDE_H - 2.6, SLIDE_W, 2.6, BG_PANEL)
                _set_shape_alpha(overlay, 82)
                _gradient_strip(slide, 0, SLIDE_H - 2.68, SLIDE_W, 0.08)
                _eyebrow(slide, 0.8, SLIDE_H - 2.3, "PRESENTATION")
                _text(slide, 0.8, SLIDE_H - 1.9, 10.5, 0.9, slide_title, 32, WHITE, bold=True)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _text(slide, 0.8, SLIDE_H - 1.0, 10.5, 0.7, subtitle, 15, ACCENT_LIGHT, spacing=1.1)

            elif hero_image and LAYOUT == "bottom_bar":
                overlay = _rect(slide, 0, 0, SLIDE_W, 2.6, BG_PANEL)
                _set_shape_alpha(overlay, 82)
                _gradient_strip(slide, 0, 2.58, SLIDE_W, 0.08)
                _eyebrow(slide, 1.0, 0.4, "PRESENTATION")
                _text(slide, 1.0, 0.8, 10.5, 0.9, slide_title, 32, WHITE, bold=True)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _text(slide, 1.0, 1.7, 10.5, 0.8, subtitle, 15, ACCENT_LIGHT, spacing=1.1)
                _text(slide, 1.0, SLIDE_H - 0.6, 8.0, 0.4, "Generated by S.T.E.W Agent", 11, WHITE)

            elif hero_image:
                # diagonal / minimal — bottom two-thirds panel
                overlay = _rect(slide, 0, SLIDE_H - 3.4, SLIDE_W, 3.4, BG_PANEL)
                _set_shape_alpha(overlay, 82)
                _rect(slide, 0, SLIDE_H - 3.42, SLIDE_W, 0.06, ACCENT)
                _eyebrow(slide, 1.0, SLIDE_H - 3.05, "PRESENTATION")
                _text(slide, 0.95, SLIDE_H - 2.6, 11.0, 1.2, slide_title, 34, WHITE, bold=True, spacing=1.0)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _rect(slide, 1.0, SLIDE_H - 1.3, 0.7, 0.045, ACCENT)
                _text(slide, 0.95, SLIDE_H - 1.1, 10.5, 0.8, subtitle, 15, ACCENT_LIGHT, spacing=1.1)

            elif LAYOUT == "side_panel":
                _rect(slide, 0, 0, 4.6, SLIDE_H, BG_PANEL)
                _rect(slide, 4.6, 0, 0.05, SLIDE_H, ACCENT)
                # Gradient strip at bottom of panel
                _gradient_strip(slide, 0, SLIDE_H - 0.5, 4.6, 0.5)
                _text(slide, 0.5, 0.6, 3.6, 2.0, "01", 90, ACCENT_DIM, bold=True)
                _eyebrow(slide, 0.5, 5.6, "PRESENTATION")
                _text(slide, 0.5, 6.0, 3.6, 0.6, _date_str(), 13, MUTED)
                # Decorative circles on right side
                _oval(slide, SLIDE_W - 2.0, -0.8, 2.5, ACCENT_DIM)
                _oval(slide, SLIDE_W - 0.8, SLIDE_H - 2.0, 1.8, ACCENT_DIM)
                _eyebrow(slide, 5.1, 2.15, "OVERVIEW")
                _text(slide, 5.05, 2.6, 7.7, 1.8, slide_title, 42, WHITE, bold=True, spacing=1.0)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _rect(slide, 5.1, 4.15, 0.7, 0.045, ACCENT)
                _text(slide, 5.05, 4.35, 7.5, 0.9, subtitle, 17, ACCENT_LIGHT, spacing=1.15)
                _text(slide, 5.05, SLIDE_H - 0.9, 7.5, 0.5, "Generated by S.T.E.W Agent", 11, MUTED)

            elif LAYOUT == "top_bar":
                _rect(slide, 0, 0, SLIDE_W, 1.2, BG_PANEL)
                _gradient_strip(slide, 0, 1.2, SLIDE_W, 0.08)
                _text(slide, 0.8, 0.3, 10.0, 0.7, slide_title, 36, WHITE, bold=True)
                _eyebrow(slide, 0.8, 2.0, "PRESENTATION")
                _rect(slide, 0.85, 2.5, 1.0, 0.05, ACCENT)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _text(slide, 0.85, 2.8, 11.0, 1.0, subtitle, 20, ACCENT_LIGHT, spacing=1.15)
                _decorative_shapes(slide, variant=2)
                _text(slide, 0.85, SLIDE_H - 0.9, 7.5, 0.5, "Generated by S.T.E.W Agent", 11, MUTED)

            elif LAYOUT == "bottom_bar":
                _rect(slide, 0, SLIDE_H - 1.5, SLIDE_W, 1.5, BG_PANEL)
                _gradient_strip(slide, 0, SLIDE_H - 1.58, SLIDE_W, 0.08)
                _text(slide, 1.0, 1.8, 11.0, 2.0, slide_title, 44, WHITE, bold=True, spacing=1.0)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _rect(slide, 1.05, 4.2, 0.7, 0.045, ACCENT)
                _text(slide, 1.0, 4.4, 11.0, 0.8, subtitle, 18, ACCENT_LIGHT, spacing=1.15)
                _text(slide, 1.0, SLIDE_H - 1.0, 10.0, 0.5, "Generated by S.T.E.W Agent", 11, MUTED)
                _eyebrow(slide, 1.0, 1.2, "PRESENTATION")
                _decorative_shapes(slide, variant=3)

            elif LAYOUT == "diagonal":
                _rect(slide, 0, 0, SLIDE_W, 0.09, ACCENT)
                _shape(slide, -1, 4.5, 6, 4, ACCENT_DIM, MSO_SHAPE.PARALLELOGRAM)
                _tri(slide, SLIDE_W - 2.0, 0, 2.0, 2.0, ACCENT_DIM)
                _text(slide, 0.5, 0.6, 3.6, 1.5, "01", 80, ACCENT_DIM, bold=True)
                _eyebrow(slide, 1.0, 1.8, "PRESENTATION")
                _text(slide, 0.95, 2.3, 11.0, 1.8, slide_title, 42, WHITE, bold=True, spacing=1.0)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _rect(slide, 1.0, 4.15, 0.7, 0.045, ACCENT)
                _text(slide, 0.95, 4.4, 10.5, 0.9, subtitle, 17, ACCENT_LIGHT, spacing=1.15)
                _text(slide, 0.95, SLIDE_H - 0.9, 7.5, 0.5, "Generated by S.T.E.W Agent", 11, MUTED)

            else:  # minimal
                _rect(slide, 0, 0, SLIDE_W, 0.09, ACCENT)
                _gradient_strip(slide, 0, 0.09, 3.0, 0.09)
                _oval(slide, SLIDE_W - 2.0, -0.8, 2.5, ACCENT_DIM)
                _eyebrow(slide, 1.0, 1.8, "PRESENTATION")
                _text(slide, 0.95, 2.3, 11.0, 2.0, slide_title, 44, WHITE, bold=True, spacing=1.0)
                subtitle = bullets[0] if bullets and bullets[0].strip() != slide_title.strip() else title
                _rect(slide, 1.0, 4.5, 0.7, 0.045, ACCENT)
                _text(slide, 0.95, 4.75, 10.5, 0.9, subtitle, 17, ACCENT_LIGHT, spacing=1.15)
                _text(slide, 0.95, SLIDE_H - 0.9, 7.5, 0.5, "Generated by S.T.E.W Agent", 11, MUTED)

        elif i == len(slides) - 1 and total_slides > 3:
            # ══ CLOSING SLIDE — special layout with call to action, hero image if available ══
            if closing_image:
                import io as _io
                try:
                    slide.shapes.add_picture(_io.BytesIO(closing_image), Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
                    overlay = _rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
                    _set_shape_alpha(overlay, 78)
                except Exception:
                    closing_image = None

            _rect(slide, 0, 0, SLIDE_W, 0.09, ACCENT)
            if not closing_image:
                _decorative_shapes(slide, variant=i % 4)
            _eyebrow(slide, 1.0, 1.0, "THANK YOU")
            _text(slide, 0.95, 1.5, 11.0, 1.5, slide_title, 38, WHITE, bold=True, spacing=1.0)
            _rect(slide, 1.0, 3.2, 1.0, 0.05, ACCENT)
            if bullets:
                _bullets(slide, 1.0, 3.6, 11.2, bullets, font_size=18, line_gap=0.72)
            # CTA block at bottom
            _rrect(slide, 1.0, SLIDE_H - 1.8, 5.0, 1.0, BG_PANEL)
            _rect(slide, 1.0, SLIDE_H - 1.8, 0.08, 1.0, ACCENT)
            _text(slide, 1.3, SLIDE_H - 1.6, 4.5, 0.4, "Generated by S.T.E.W Agent", 13, ACCENT_LIGHT, bold=True)
            _text(slide, 1.3, SLIDE_H - 1.1, 4.5, 0.3, "stew-agent.onrender.com", 12, MUTED)
            _page_footer(slide, i)

        else:
            # ══ CONTENT SLIDE — alternating layouts for variety ══
            content_variant = i % 3  # alternate between 3 content layouts

            if content_variant == 0:
                # Layout A: Standard with decorative shapes
                _rect(slide, 0, 0, SLIDE_W, 0.09, ACCENT)
                _decorative_shapes(slide, variant=i % 4)
                _eyebrow(slide, 0.9, 0.45, f"SECTION {i:02d}")
                _text(slide, 0.85, 0.85, 11.0, 0.95, slide_title, 30, WHITE, bold=True)
                _rect(slide, 0.9, 1.75, 1.0, 0.05, ACCENT)
                _text(slide, 11.3, 0.35, 1.6, 1.2, f"{i:02d}", 46, ACCENT_DIM, bold=True, align=PP_ALIGN.RIGHT)
                if bullets:
                    _bullets(slide, 1.0, 2.25, 11.2, bullets, font_size=17.5, line_gap=0.68)
                _page_footer(slide, i)

            elif content_variant == 1:
                # Layout B: Two-column with accent panel on right
                _rect(slide, 0, 0, SLIDE_W, 0.09, ACCENT)
                # Right accent panel
                _rect(slide, SLIDE_W - 3.5, 0, 3.5, SLIDE_H, BG_PANEL)
                _rect(slide, SLIDE_W - 3.55, 0, 0.05, SLIDE_H, ACCENT)
                # Large number in panel
                _text(slide, SLIDE_W - 3.2, 0.5, 3.0, 2.5, f"{i:02d}", 72, ACCENT_DIM, bold=True, align=PP_ALIGN.CENTER)
                _eyebrow(slide, SLIDE_W - 3.2, 3.0, "SECTION")
                # Left content
                _eyebrow(slide, 0.9, 0.45, f"SECTION {i:02d}")
                _text(slide, 0.85, 0.85, 9.0, 0.95, slide_title, 28, WHITE, bold=True)
                _rect(slide, 0.9, 1.75, 1.0, 0.05, ACCENT)
                if bullets:
                    _bullets(slide, 1.0, 2.25, 8.8, bullets, font_size=17, line_gap=0.68)
                _page_footer(slide, i)

            else:
                # Layout C: Card-based layout — bullets as cards
                _rect(slide, 0, 0, SLIDE_W, 0.09, ACCENT)
                _eyebrow(slide, 0.9, 0.45, f"SECTION {i:02d}")
                _text(slide, 0.85, 0.85, 11.0, 0.95, slide_title, 30, WHITE, bold=True)
                _rect(slide, 0.9, 1.75, 1.0, 0.05, ACCENT)
                _text(slide, 11.3, 0.35, 1.6, 1.2, f"{i:02d}", 46, ACCENT_DIM, bold=True, align=PP_ALIGN.RIGHT)

                if bullets and len(bullets) <= 3:
                    # Render as cards
                    card_w = 3.6
                    card_gap = 0.3
                    start_x = (SLIDE_W - (len(bullets) * card_w + (len(bullets) - 1) * card_gap)) / 2
                    for j, bullet in enumerate(bullets):
                        clean = _sanitize(bullet.strip().lstrip("-").lstrip("*").lstrip("\u2022").strip())
                        if not clean:
                            continue
                        cx = start_x + j * (card_w + card_gap)
                        _card(slide, cx, 2.5, card_w, 3.5, f"Point {j+1}", clean)
                elif bullets:
                    # Too many bullets for cards — use standard layout
                    _bullets(slide, 1.0, 2.25, 11.2, bullets, font_size=17.5, line_gap=0.68)

                # Decorative accent
                _oval(slide, -0.5, SLIDE_H - 2.0, 1.5, ACCENT_DIM)
                _page_footer(slide, i)


def _sanitize(text):
    """Clean text for PPTX rendering — strips ALL markdown markers."""
    if not text:
        return ""
    import re as _re
    # Unicode replacements
    replacements = {
        "\u2248": "~", "\u00d7": "x", "\u2212": "-", "\u2013": "-", "\u2014": "--",
        "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
        "\u2026": "...", "\u00a0": " ",
    }
    for old_char, new_char in replacements.items():
        text = text.replace(old_char, new_char)
    # Strip heading markers: ## Heading -> Heading
    text = _re.sub(r'^#{1,6}\s+', '', text, flags=_re.MULTILINE)
    # Strip bold: **text** or __text__ -> text
    text = _re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = _re.sub(r'__(.+?)__', r'\1', text)
    # Strip italic: *text* or _text_ -> text
    text = _re.sub(r'(?<!\w)\*([^*\n]+?)\*(?!\w)', r'\1', text)
    text = _re.sub(r'(?<!\w)_([^_\n]+?)_(?!\w)', r'\1', text)
    # Strip inline code backticks
    text = _re.sub(r'`([^`\n]+?)`', r'\1', text)
    # Strip horizontal rules
    text = _re.sub(r'^[\s]*[-_]{3,}[\s]*$', '', text, flags=_re.MULTILINE)
    # Strip blockquotes
    text = _re.sub(r'^>\s+', '', text, flags=_re.MULTILINE)
    return text.strip()


def _date_str():
    from datetime import datetime
    return datetime.utcnow().strftime("%B %Y")
