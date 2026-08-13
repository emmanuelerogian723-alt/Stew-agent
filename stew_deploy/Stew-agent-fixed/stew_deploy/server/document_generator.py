"""
S.T.E.W Document Generator — REAL binary file generation.
Supports: PDF, DOCX, XLSX, PPTX, HTML
All files returned as base64-encoded strings.
"""
import base64
import io
import json
import logging
from datetime import datetime
from typing import Any, Optional

from fastapi import HTTPException

logger = logging.getLogger(__name__)


def _to_base64(buf: io.BytesIO) -> str:
    buf.seek(0)
    return base64.b64encode(buf.read()).decode("utf-8")


# Common unicode chars that break in reportlab's default Helvetica font
_UNICODE_REPLACEMENTS = {
    "\u2248": "~",     # ≈
    "\u00d7": "x",     # ×
    "\u2212": "-",     # − (unicode minus)
    "\u2013": "-",     # – (en dash)
    "\u2014": "--",    # — (em dash)
    "\u2018": "'", "\u2019": "'",  # smart quotes
    "\u201c": '"', "\u201d": '"',
    "\u2192": "->",    # →
    "\u03bb": "lambda", "\u039b": "Lambda",
    "\u03c1": "rho", "\u03b4": "delta",
    "\u03bc": "u", "\u00b5": "u",  # micro/mu
    "\u2103": " deg C",
    "\u00b1": "+/-",   # ±
    "\u2265": ">=", "\u2264": "<=",
    "\u2260": "!=",
    "\u221e": "infinity",
    "\u00b0": " deg",
    "\u25a0": "",      # black square (the ■ artifact itself)
    "\ufffd": "",      # replacement character
}

# Superscript/subscript digit maps -> plain ASCII (e.g. 10⁻³² -> 10^-32)
_SUPERSCRIPT_MAP = str.maketrans({
    "\u2070": "0", "\u00b9": "1", "\u00b2": "2", "\u00b3": "3",
    "\u2074": "4", "\u2075": "5", "\u2076": "6", "\u2077": "7",
    "\u2078": "8", "\u2079": "9", "\u207b": "-",
})
_SUBSCRIPT_MAP = str.maketrans({
    "\u2080": "0", "\u2081": "1", "\u2082": "2", "\u2083": "3",
    "\u2084": "4", "\u2085": "5", "\u2086": "6", "\u2087": "7",
    "\u2088": "8", "\u2089": "9", "\u208b": "-",
})


def _sanitize_text(text: str) -> str:
    """Strip/replace unicode chars that render as black boxes in default PDF/DOCX fonts."""
    if not text:
        return text
    text = text.translate(_SUPERSCRIPT_MAP)
    text = text.translate(_SUBSCRIPT_MAP)
    for uni, replacement in _UNICODE_REPLACEMENTS.items():
        text = text.replace(uni, replacement)
    # Strip any remaining non-ASCII characters (final safety net)
    text = text.encode("ascii", "ignore").decode("ascii")
    return text


# ── PDF ───────────────────────────────────────────────────────────────────────

def _parse_md_table_row(line: str) -> list:
    """Parse a markdown table row '| a | b | c |' into ['a','b','c']."""
    cells = [c.strip() for c in line.strip().strip("|").split("|")]
    return [_sanitize_text(c.replace("**", "").replace("*", "")) for c in cells]


def _is_table_separator(line: str) -> bool:
    """Detect '|---|---|---|' style separator rows."""
    stripped = line.strip()
    if not stripped.startswith("|"):
        return False
    inner = stripped.strip("|")
    return all(c in "-:| " for c in inner) and "-" in inner


def generate_pdf(content: str, title: str = "Document") -> dict:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER, TA_LEFT

        buf = io.BytesIO()
        doc = SimpleDocTemplate(
            buf,
            pagesize=A4,
            rightMargin=2 * cm,
            leftMargin=2 * cm,
            topMargin=2 * cm,
            bottomMargin=2 * cm,
            title=_sanitize_text(title),
            author="S.T.E.W Agent",
        )

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Title"],
            fontSize=18,
            spaceAfter=12,
            alignment=TA_CENTER,
        )
        body_style = ParagraphStyle(
            "CustomBody",
            parent=styles["Normal"],
            fontSize=11,
            leading=16,
            spaceAfter=8,
        )
        table_style = TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563EB")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#EFF6FF")]),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ])

        story = []
        story.append(Paragraph(_sanitize_text(title), title_style))
        story.append(Spacer(1, 0.5 * cm))

        lines = content.split("\n")
        i = 0
        n = len(lines)
        while i < n:
            line = lines[i]
            stripped = line.strip()

            # Detect start of a markdown table: a row with pipes, followed by a separator row
            if stripped.startswith("|") and i + 1 < n and _is_table_separator(lines[i + 1]):
                header_row = _parse_md_table_row(stripped)
                table_data = [header_row]
                j = i + 2
                while j < n and lines[j].strip().startswith("|"):
                    table_data.append(_parse_md_table_row(lines[j]))
                    j += 1
                # Wrap cell text in Paragraphs so long text wraps instead of overflowing
                cell_style = ParagraphStyle("Cell", parent=styles["Normal"], fontSize=9, leading=11)
                wrapped = [[Paragraph(str(cell), cell_style) for cell in row] for row in table_data]
                col_count = len(header_row) or 1
                avail_width = 17 * cm
                col_width = avail_width / col_count
                t = Table(wrapped, colWidths=[col_width] * col_count, repeatRows=1)
                t.setStyle(table_style)
                story.append(t)
                story.append(Spacer(1, 0.4 * cm))
                i = j
                continue

            if not stripped:
                story.append(Spacer(1, 0.3 * cm))
            elif stripped.startswith("## "):
                clean_h = _sanitize_text(stripped[3:].replace("**", "").replace("*", ""))
                story.append(Paragraph(clean_h, styles["Heading2"]))
            elif stripped.startswith("# "):
                clean_h = _sanitize_text(stripped[2:].replace("**", "").replace("*", ""))
                story.append(Paragraph(clean_h, styles["Heading1"]))
            elif stripped.startswith("- ") or stripped.startswith("* ") or stripped.startswith("• "):
                bullet_text = _sanitize_text(stripped.lstrip("-*• ").replace("**", "").replace("*", ""))
                story.append(Paragraph(f"• {bullet_text}", body_style))
            else:
                # Strip **bold** and *italic* markers, sanitize unicode, escape HTML chars
                safe = _sanitize_text(stripped.replace("**", "").replace("*", ""))
                safe = safe.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(safe, body_style))
            i += 1

        # Footer with timestamp
        story.append(Spacer(1, 1 * cm))
        footer_style = ParagraphStyle("Footer", parent=styles["Normal"], fontSize=8,
                                      textColor=colors.grey, alignment=TA_CENTER)
        story.append(Paragraph(
            f"Generated by S.T.E.W Agent - {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}",
            footer_style,
        ))

        doc.build(story)
        filename = f"{title.replace(' ', '_')}_{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.pdf"
        return {
            "file": _to_base64(buf),
            "filename": filename,
            "mime_type": "application/pdf",
            "success": True,
        }
    except ImportError:
        raise HTTPException(500, "reportlab not installed — add it to requirements.txt")
    except Exception as e:
        logger.error(f"PDF generation error: {e}")
        raise HTTPException(500, f"PDF generation failed: {e}")


# ── DOCX ──────────────────────────────────────────────────────────────────────

def generate_docx(content: str, title: str = "Document") -> dict:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()

        # Document title
        title_para = doc.add_heading(_sanitize_text(title), level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph()  # spacer

        def _shade_cell(cell, color_hex):
            shading = OxmlElement("w:shd")
            shading.set(qn("w:fill"), color_hex)
            cell._tc.get_or_add_tcPr().append(shading)

        lines = content.split("\n")
        i = 0
        n = len(lines)
        while i < n:
            line = lines[i]
            stripped = line.strip()

            # Detect markdown table
            if stripped.startswith("|") and i + 1 < n and _is_table_separator(lines[i + 1]):
                header_row = _parse_md_table_row(stripped)
                table_rows = []
                j = i + 2
                while j < n and lines[j].strip().startswith("|"):
                    table_rows.append(_parse_md_table_row(lines[j]))
                    j += 1

                col_count = len(header_row) or 1
                table = doc.add_table(rows=1, cols=col_count)
                table.style = "Table Grid"
                hdr_cells = table.rows[0].cells
                for idx, htext in enumerate(header_row):
                    hdr_cells[idx].text = htext
                    _shade_cell(hdr_cells[idx], "2563EB")
                    for p in hdr_cells[idx].paragraphs:
                        for run in p.runs:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                for row_idx, row_data in enumerate(table_rows):
                    row_cells = table.add_row().cells
                    for idx in range(col_count):
                        val = row_data[idx] if idx < len(row_data) else ""
                        row_cells[idx].text = val
                        if row_idx % 2 == 1:
                            _shade_cell(row_cells[idx], "EFF6FF")
                doc.add_paragraph()  # spacer after table
                i = j
                continue

            if not stripped:
                doc.add_paragraph()
            elif stripped.startswith("## "):
                doc.add_heading(_sanitize_text(stripped[3:].replace("**", "").replace("*", "")), level=2)
            elif stripped.startswith("### "):
                doc.add_heading(_sanitize_text(stripped[4:].replace("**", "").replace("*", "")), level=3)
            elif stripped.startswith("# "):
                doc.add_heading(_sanitize_text(stripped[2:].replace("**", "").replace("*", "")), level=1)
            elif stripped.startswith("- ") or stripped.startswith("* ") or stripped.startswith("• "):
                clean_bullet = _sanitize_text(stripped.lstrip("-*• ").replace("**", "").replace("*", ""))
                doc.add_paragraph(clean_bullet, style="List Bullet")
            elif stripped.startswith("1. ") or (len(stripped) > 2 and stripped[0].isdigit() and stripped[1] == "."):
                doc.add_paragraph(_sanitize_text(stripped), style="List Number")
            else:
                doc.add_paragraph(_sanitize_text(stripped.replace("**", "").replace("*", "")))
            i += 1

        # Footer paragraph
        footer = doc.add_paragraph()
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = footer.add_run(
            f"Generated by S.T.E.W Agent - {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
        )
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        buf = io.BytesIO()
        doc.save(buf)
        filename = f"{title.replace(' ', '_')}_{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.docx"
        return {
            "file": _to_base64(buf),
            "filename": filename,
            "mime_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "success": True,
        }
    except ImportError:
        raise HTTPException(500, "python-docx not installed")
    except Exception as e:
        logger.error(f"DOCX generation error: {e}")
        raise HTTPException(500, f"DOCX generation failed: {e}")


# ── XLSX ──────────────────────────────────────────────────────────────────────

def generate_xlsx(
    data: list[dict],
    sheet_name: str = "Sheet1",
    title: str = "Spreadsheet",
) -> dict:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name[:31]  # Excel sheet name limit

        if not data:
            data = [{"No data": "provided"}]

        headers = list(data[0].keys())

        # Header row styling
        header_fill = PatternFill(start_color="2563EB", end_color="2563EB", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=11)
        border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = border

        # Data rows
        alt_fill = PatternFill(start_color="EFF6FF", end_color="EFF6FF", fill_type="solid")
        for row_idx, row in enumerate(data, 2):
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=row.get(header, ""))
                cell.border = border
                if row_idx % 2 == 0:
                    cell.fill = alt_fill

        # Auto-fit columns
        for col_idx, header in enumerate(headers, 1):
            max_len = max(
                len(str(header)),
                *[len(str(row.get(header, ""))) for row in data],
            )
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 4, 50)

        ws.freeze_panes = "A2"

        buf = io.BytesIO()
        wb.save(buf)
        filename = f"{title.replace(' ', '_')}_{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.xlsx"
        return {
            "file": _to_base64(buf),
            "filename": filename,
            "mime_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "success": True,
        }
    except ImportError:
        raise HTTPException(500, "openpyxl not installed")
    except Exception as e:
        logger.error(f"XLSX generation error: {e}")
        raise HTTPException(500, f"XLSX generation failed: {e}")


# ── PPTX ──────────────────────────────────────────────────────────────────────

def generate_pptx(slides: list[dict], title: str = "Presentation") -> dict:
    """Generate an advanced, premium PPTX — dark navy + gold, numbered slides, shape-based bullet markers, eyebrow labels."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn

        prs = Presentation()
        SLIDE_W = 13.333
        SLIDE_H = 7.5
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)

        # ── Premium dark + gold palette ──
        BG_DARK = RGBColor(0x0B, 0x0F, 0x1A)       # near-black navy
        BG_PANEL = RGBColor(0x14, 0x1B, 0x2E)      # side panel shade
        ACCENT = RGBColor(0xF5, 0x9E, 0x0B)        # amber gold
        ACCENT_LIGHT = RGBColor(0xFD, 0xD8, 0x6E)  # light gold
        ACCENT_DIM = RGBColor(0x7A, 0x5A, 0x1E)    # muted gold (faint numbers)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT = RGBColor(0xD7, 0xDE, 0xE8)
        MUTED = RGBColor(0x7C, 0x8A, 0xA0)
        DIVIDER = RGBColor(0x26, 0x2F, 0x45)

        BLANK = prs.slide_layouts[6]
        total_slides = len(slides)

        def _add_bg(slide, color):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def _rect(slide, left, top, width, height, color, line=False):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            if not line:
                shape.line.fill.background()
            else:
                shape.line.color.rgb = color
                shape.line.width = Pt(0.5)
            shape.shadow.inherit = False
            return shape

        def _oval(slide, left, top, size, color):
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            shape.shadow.inherit = False
            return shape

        def _text(slide, left, top, width, height, text, size, color,
                  bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  font="Segoe UI", spacing=1.0, letter_spacing=None):
            box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = anchor
            para = tf.paragraphs[0]
            para.text = _sanitize_text(text)
            para.alignment = align
            para.line_spacing = spacing
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font
            return tf

        def _eyebrow(slide, left, top, text):
            """Small uppercase gold label tag above a title."""
            _text(slide, left, top, 8.0, 0.4, text.upper(), 12.5, ACCENT,
                  bold=True, font="Segoe UI Semibold")

        def _page_footer(slide, index):
            """Bottom-right slide counter + thin gold rule."""
            _rect(slide, 0, SLIDE_H - 0.06, SLIDE_W, 0.06, ACCENT)
            _text(slide, SLIDE_W - 2.2, SLIDE_H - 0.55, 1.9, 0.4,
                  f"{index+1:02d}  /  {total_slides:02d}", 11, MUTED, align=PP_ALIGN.RIGHT)
            _text(slide, 0.6, SLIDE_H - 0.55, 4.0, 0.4, "S.T.E.W", 11, MUTED, bold=True)

        def _bulleted_block(slide, left, top, width, bullets, font_size=17.5,
                            line_gap=0.62, color=LIGHT, marker_color=ACCENT):
            """Render bullets as gold square markers + wrapped text, evenly spaced."""
            y = top
            for bullet in bullets:
                clean = _sanitize_text(bullet.strip().lstrip("-").lstrip("*").lstrip("\u2022").strip())
                if not clean:
                    continue
                # gold marker square
                _rect(slide, left, y + 0.12, 0.14, 0.14, marker_color)
                # bullet text
                box = slide.shapes.add_textbox(Inches(left + 0.35), Inches(y), Inches(width - 0.35), Inches(line_gap + 0.3))
                tf = box.text_frame
                tf.word_wrap = True
                para = tf.paragraphs[0]
                para.text = clean
                para.line_spacing = 1.05
                run = para.runs[0] if para.runs else para.add_run()
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.name = "Segoe UI"
                y += line_gap
            return y

        for i, slide_data in enumerate(slides):
            slide_title = slide_data.get("title", f"Slide {i+1}")
            slide_content = slide_data.get("content", "")
            bullets = [b.strip() for b in slide_content.split("\n") if b.strip()]

            slide = prs.slides.add_slide(BLANK)
            _add_bg(slide, BG_DARK)

            if i == 0:
                # ══ TITLE SLIDE — big cover with side panel ══
                _rect(slide, 0, 0, 4.6, SLIDE_H, BG_PANEL)          # left dark panel
                _rect(slide, 4.6, 0, 0.05, SLIDE_H, ACCENT)          # gold seam

                # Giant faint slide number on the panel
                _text(slide, 0.5, 0.6, 3.6, 2.0, "01", 90, ACCENT_DIM, bold=True)
                _eyebrow(slide, 0.5, 5.6, "PRESENTATION")
                _text(slide, 0.5, 6.0, 3.6, 0.6, datetime.utcnow().strftime("%B %Y"), 13, MUTED)

                # Main title + subtitle on the right
                _eyebrow(slide, 5.1, 2.15, "OVERVIEW")
                _text(slide, 5.05, 2.6, 7.7, 1.8, slide_title, 42, WHITE, bold=True, spacing=1.0)
                subtitle = bullets[0] if bullets else title
                if subtitle.strip() == slide_title.strip():
                    subtitle = title
                _rect(slide, 5.1, 4.15, 0.7, 0.045, ACCENT)
                _text(slide, 5.05, 4.35, 7.5, 0.9, subtitle, 17, ACCENT_LIGHT, spacing=1.15)
                _text(slide, 5.05, SLIDE_H - 0.9, 7.5, 0.5, "Generated by S.T.E.W Agent", 11, MUTED)

            else:
                # ══ CONTENT SLIDE ══
                # thin top accent + eyebrow section tag
                _rect(slide, 0, 0, SLIDE_W, 0.09, ACCENT)
                _eyebrow(slide, 0.9, 0.45, f"SECTION {i:02d}")
                _text(slide, 0.85, 0.85, 11.0, 0.95, slide_title, 30, WHITE, bold=True)
                _rect(slide, 0.9, 1.75, 1.0, 0.05, ACCENT)

                # faint big number top-right (design accent)
                _text(slide, 11.3, 0.35, 1.6, 1.2, f"{i:02d}", 46, ACCENT_DIM, bold=True, align=PP_ALIGN.RIGHT)

                if bullets:
                    _bulleted_block(slide, 1.0, 2.25, 11.2, bullets, font_size=17.5, line_gap=0.68)

                _page_footer(slide, i)

        buf = io.BytesIO()
        prs.save(buf)
        clean_title = title.replace(" ", "_").replace("/", "_")[:40]
        filename = f"{clean_title}_{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.pptx"
        return {
            "file": _to_base64(buf),
            "filename": filename,
            "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "success": True,
        }
    except ImportError:
        raise HTTPException(500, "python-pptx not installed")
    except Exception as e:
        logger.error(f"PPTX generation error: {e}")
        raise HTTPException(500, f"PPTX generation failed: {e}")

# ── HTML ──────────────────────────────────────────────────────────────────────

def generate_html(content: str, title: str = "Report") -> dict:
    import html as html_module

    # Convert markdown-ish content to HTML
    lines = content.split("\n")
    body_html = []
    in_ul = False

    for line in lines:
        stripped = line.strip()
        if not stripped:
            if in_ul:
                body_html.append("</ul>")
                in_ul = False
            body_html.append("<br>")
        elif stripped.startswith("## "):
            if in_ul:
                body_html.append("</ul>")
                in_ul = False
            body_html.append(f"<h2>{html_module.escape(stripped[3:])}</h2>")
        elif stripped.startswith("# "):
            if in_ul:
                body_html.append("</ul>")
                in_ul = False
            body_html.append(f"<h1>{html_module.escape(stripped[2:])}</h1>")
        elif stripped.startswith("### "):
            if in_ul:
                body_html.append("</ul>")
                in_ul = False
            body_html.append(f"<h3>{html_module.escape(stripped[4:])}</h3>")
        elif stripped.startswith("- ") or stripped.startswith("* "):
            if not in_ul:
                body_html.append("<ul>")
                in_ul = True
            body_html.append(f"<li>{html_module.escape(stripped[2:])}</li>")
        else:
            if in_ul:
                body_html.append("</ul>")
                in_ul = False
            body_html.append(f"<p>{html_module.escape(stripped)}</p>")

    if in_ul:
        body_html.append("</ul>")

    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{html_module.escape(title)}</title>
  <style>
    *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{
      font-family: 'Segoe UI', system-ui, -apple-system, sans-serif;
      background: #0f172a;
      color: #e2e8f0;
      line-height: 1.7;
      padding: 2rem;
    }}
    .container {{ max-width: 900px; margin: 0 auto; }}
    header {{
      background: linear-gradient(135deg, #1e293b, #0f172a);
      border-left: 4px solid #38bdf8;
      padding: 1.5rem 2rem;
      margin-bottom: 2rem;
      border-radius: 0 8px 8px 0;
    }}
    header h1 {{ color: #38bdf8; font-size: 2rem; }}
    header small {{ color: #94a3b8; font-size: 0.85rem; }}
    .content {{ background: #1e293b; padding: 2rem; border-radius: 8px; }}
    h1, h2, h3 {{ color: #38bdf8; margin: 1.5rem 0 0.75rem; }}
    h1 {{ font-size: 1.8rem; }}
    h2 {{ font-size: 1.4rem; border-bottom: 1px solid #334155; padding-bottom: 0.4rem; }}
    h3 {{ font-size: 1.1rem; color: #7dd3fc; }}
    p {{ margin-bottom: 0.8rem; color: #cbd5e1; }}
    ul {{ padding-left: 1.5rem; margin-bottom: 0.8rem; }}
    li {{ margin-bottom: 0.4rem; color: #cbd5e1; }}
    footer {{
      text-align: center;
      margin-top: 2rem;
      color: #64748b;
      font-size: 0.8rem;
    }}
    @media (max-width: 600px) {{ body {{ padding: 1rem; }} }}
  </style>
</head>
<body>
  <div class="container">
    <header>
      <h1>{html_module.escape(title)}</h1>
      <small>Generated by S.T.E.W Agent • {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}</small>
    </header>
    <div class="content">
      {''.join(body_html)}
    </div>
    <footer>S.T.E.W — Structured Task Execution Workflow</footer>
  </div>
</body>
</html>"""

    buf = io.BytesIO(html_content.encode("utf-8"))
    filename = f"{title.replace(' ', '_')}_{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.html"
    return {
        "file": _to_base64(buf),
        "filename": filename,
        "mime_type": "text/html",
        "success": True,
    }
